import xlrd
from pptx import Presentation
import re
import string

class CertificatorGeneratorTabajara():
    
    def __init__(self) -> None:
        loc = ("info_certificados.xls")
        wb = xlrd.open_workbook(loc)
        self.sheet = wb.sheet_by_name('Informacoes')

    def reading_information(self):       

        self.sheet.cell_value(0, 0)
        self.values = []
        for i in range (self.sheet.nrows):
            if i !=0:
                titulo = self.sheet.cell_value(i, 0)
                apresentador = self.sheet.cell_value(i, 1)
                dados = [titulo, apresentador]
                self.values.append (dados)
    
    def create_cerification(self):

        for value in self.values:
            prs = Presentation(".\\templates\\template_1.pptx") 
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    count = 0
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "XXXXXXXXXXXXX" in run.text:
                                run.text = run.text.replace ('XXXXXXXXXXXXX', (value[count]).strip())  
                                count = count +1
                                
            chars = re.escape(string.punctuation)
            valueX = re.sub(r'['+chars+']', '',value[0].strip())
            file_name = "Certificado-"+valueX.replace(' ','_')
            
            prs.save('.\\certificados\\'+file_name+'.pptx') 

    def do (self):
        self.reading_information()
        self.create_cerification()
        

certificator_generator_tabajara = CertificatorGeneratorTabajara()
certificator_generator_tabajara.do()